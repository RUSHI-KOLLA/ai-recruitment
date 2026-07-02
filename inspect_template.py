"""Inspect the PPTX template to understand slide layouts and placeholders."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import sys

prs = Presentation(sys.argv[1])

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Number of existing slides: {len(prs.slides)}")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")
print("=" * 70)

# Show existing slides
for i, slide in enumerate(prs.slides):
    print(f"\n--- EXISTING SLIDE {i} ---")
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        print(f"  Shape: name='{shape.name}', type={stype}, "
              f"pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})")
        if shape.has_text_frame:
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if text:
                    font_info = ""
                    if para.runs:
                        r = para.runs[0]
                        font_info = f" [font={r.font.name}, size={r.font.size}, bold={r.font.bold}, color={r.font.color.rgb if r.font.color and r.font.color.rgb else 'N/A'}]"
                    print(f"    Para[{p_idx}]: '{text[:100]}'{font_info}")
